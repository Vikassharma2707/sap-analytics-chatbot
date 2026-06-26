"""
Export Service — generates PowerPoint, PDF, and Excel exports from analytics results.
"""

from __future__ import annotations
import os
import uuid
from typing import Any
from pathlib import Path
import structlog

from app.core.config import settings

log = structlog.get_logger()

os.makedirs(settings.EXPORT_DIR, exist_ok=True)


def _export_path(fmt: str) -> str:
    return str(Path(settings.EXPORT_DIR) / f"{uuid.uuid4()}.{fmt}")


def export_excel(
    records: list[dict],
    kpis: dict[str, Any],
    title: str = "SAP Analytics Export",
) -> str:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    import pandas as pd

    path = _export_path("xlsx")
    wb = openpyxl.Workbook()

    # KPI sheet
    ws_kpi = wb.active
    ws_kpi.title = "KPIs"
    ws_kpi["A1"] = title
    ws_kpi["A1"].font = Font(size=14, bold=True)
    ws_kpi["A1"].fill = PatternFill("solid", fgColor="1F4E79")

    row = 3
    for measure, stats in kpis.items():
        ws_kpi.cell(row=row, column=1, value=measure).font = Font(bold=True)
        for col_i, (k, v) in enumerate(stats.items(), start=2):
            ws_kpi.cell(row=row, column=col_i, value=f"{k}: {round(v, 2) if isinstance(v, float) else v}")
        row += 1

    # Data sheet
    ws_data = wb.create_sheet("Data")
    if records:
        df = __import__("pandas").DataFrame(records)
        header_fill = PatternFill("solid", fgColor="2E75B6")
        for col_i, col_name in enumerate(df.columns, start=1):
            cell = ws_data.cell(row=1, column=col_i, value=col_name)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
        for row_i, row_data in enumerate(df.itertuples(index=False), start=2):
            for col_i, value in enumerate(row_data, start=1):
                ws_data.cell(row=row_i, column=col_i, value=value)
        # Auto-width
        for col in ws_data.columns:
            max_len = max(len(str(c.value or "")) for c in col)
            ws_data.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)

    wb.save(path)
    log.info("excel_exported", path=path)
    return path


def export_pptx(
    title: str,
    insights: dict[str, Any],
    kpis: dict[str, Any],
    records: list[dict],
    chart_image_path: str | None = None,
) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = _export_path("pptx")
    prs = Presentation()
    prs.slide_width = __import__("pptx.util", fromlist=["Emu"]).Emu(9144000)
    prs.slide_height = __import__("pptx.util", fromlist=["Emu"]).Emu(5143500)

    BLUE = RGBColor(0x1F, 0x4E, 0x79)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def add_text_box(slide, text, left, top, width, height, size=18, bold=False, color=None):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    # Slide 1: Title
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = BLUE
    add_text_box(slide1, title, 0.5, 1.5, 9, 1.5, size=36, bold=True, color=WHITE)
    add_text_box(slide1, "SAP S/4HANA Analytics Report", 0.5, 3.2, 9, 0.8, size=20, color=WHITE)

    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(blank_layout)
    add_text_box(slide2, "Executive Summary", 0.3, 0.2, 9, 0.6, size=24, bold=True, color=BLUE)
    summary = insights.get("summary", "")
    add_text_box(slide2, summary, 0.3, 1.0, 9, 1.2, size=16)

    row = 2.5
    for insight in insights.get("key_insights", [])[:5]:
        add_text_box(slide2, f"• {insight}", 0.5, row, 8.5, 0.5, size=14)
        row += 0.55

    # Slide 3: KPIs
    slide3 = prs.slides.add_slide(blank_layout)
    add_text_box(slide3, "Key Performance Indicators", 0.3, 0.2, 9, 0.6, size=24, bold=True, color=BLUE)
    row = 1.0
    for measure, stats in list(kpis.items())[:4]:
        add_text_box(slide3, f"{measure}: {round(stats.get('total', 0), 2):,.2f}", 0.5, row, 8, 0.5, size=18, bold=True)
        row += 0.7

    # Slide 4: Recommendations
    slide4 = prs.slides.add_slide(blank_layout)
    add_text_box(slide4, "Recommendations", 0.3, 0.2, 9, 0.6, size=24, bold=True, color=BLUE)
    row = 1.0
    for rec in insights.get("recommendations", [])[:4]:
        add_text_box(slide4, f"➤ {rec}", 0.5, row, 8.5, 0.6, size=15)
        row += 0.8

    prs.save(path)
    log.info("pptx_exported", path=path)
    return path


def export_pdf(
    title: str,
    insights: dict[str, Any],
    kpis: dict[str, Any],
    records: list[dict],
) -> str:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                     Table, TableStyle, HRFlowable)

    path = _export_path("pdf")
    doc = SimpleDocTemplate(path, pagesize=letter, topMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    story = []

    BLUE = colors.HexColor("#1F4E79")
    title_style = ParagraphStyle("title", parent=styles["Heading1"], textColor=BLUE, fontSize=20)
    h2_style = ParagraphStyle("h2", parent=styles["Heading2"], textColor=BLUE, fontSize=14)

    story.append(Paragraph(title, title_style))
    story.append(Paragraph("SAP S/4HANA Analytics Report", styles["Normal"]))
    story.append(Spacer(1, 0.3 * inch))
    story.append(HRFlowable(width="100%", color=BLUE))
    story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("Executive Summary", h2_style))
    story.append(Paragraph(insights.get("summary", ""), styles["Normal"]))
    story.append(Spacer(1, 0.2 * inch))

    story.append(Paragraph("Key Insights", h2_style))
    for insight in insights.get("key_insights", []):
        story.append(Paragraph(f"• {insight}", styles["Normal"]))
    story.append(Spacer(1, 0.2 * inch))

    if kpis:
        story.append(Paragraph("KPIs", h2_style))
        kpi_data = [["Metric", "Total", "Average", "Count"]]
        for measure, stats in list(kpis.items())[:6]:
            kpi_data.append([
                measure,
                f"{stats.get('total', 0):,.2f}",
                f"{stats.get('average', 0):,.2f}",
                str(stats.get("count", 0)),
            ])
        tbl = Table(kpi_data, colWidths=[2.5 * inch, 1.5 * inch, 1.5 * inch, 1.5 * inch])
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), BLUE),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.lightblue]),
        ]))
        story.append(tbl)

    if records:
        story.append(Spacer(1, 0.3 * inch))
        story.append(Paragraph("Data Summary (Top 20)", h2_style))
        headers = list(records[0].keys())[:6]
        table_data = [headers] + [[str(r.get(h, ""))[:30] for h in headers] for r in records[:20]]
        col_w = 7.5 * inch / len(headers)
        tbl2 = Table(table_data, colWidths=[col_w] * len(headers))
        tbl2.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), BLUE),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
        ]))
        story.append(tbl2)

    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph("Recommendations", h2_style))
    for rec in insights.get("recommendations", []):
        story.append(Paragraph(f"➤ {rec}", styles["Normal"]))

    doc.build(story)
    log.info("pdf_exported", path=path)
    return path
